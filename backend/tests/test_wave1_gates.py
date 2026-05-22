"""Integration tests for Wave 1 validation gates."""

import asyncio
import json
from pathlib import Path
from typing import Optional

import pytest
import tempfile

# Note: These tests require the extraction modules and dependencies to be installed
# Run: pip install -r backend/requirements.txt


@pytest.mark.asyncio
async def test_health_endpoint():
    """Test 1: GET /health returns 200 with correct response."""
    try:
        from backend.main import app
        from fastapi.testclient import TestClient
        
        client = TestClient(app)
        response = client.get("/health")
        
        assert response.status_code == 200, f"Expected 200, got {response.status_code}"
        
        data = response.json()
        assert data["status"] == "ok"
        assert "version" in data
        assert "ollama_available" in data
        
        print("✅ Health endpoint test passed")
        
    except ImportError as e:
        pytest.skip(f"Dependencies not installed: {e}")


@pytest.mark.asyncio
async def test_pdf_extraction():
    """Test 2: PyMuPDF extracts readable text from sample PDF."""
    try:
        from backend.extraction.pdf import extract_text_from_pdf
        
        # Create a test PDF
        import fitz
        
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            pdf_path = f.name
        
        try:
            # Create minimal PDF
            doc = fitz.open()
            page = doc.new_page()
            page.insert_text((50, 50), "Introduction to Computer Science\nCMSC 11\nWeek 1")
            doc.save(pdf_path)
            doc.close()
            
            # Extract text
            text = await extract_text_from_pdf(pdf_path)
            
            assert text, "No text extracted from PDF"
            assert "Computer Science" in text or "CMSC" in text or "Week" in text
            assert len(text) > 20, "Extracted text too short"
            
            print("✅ PDF extraction test passed")
            
        finally:
            Path(pdf_path).unlink(missing_ok=True)
            
    except ImportError as e:
        pytest.skip(f"Dependencies not installed: {e}")


@pytest.mark.asyncio
async def test_docx_extraction():
    """Test 3: python-docx extracts text from DOCX files."""
    try:
        from backend.extraction.docx import extract_text_from_docx
        from docx import Document
        
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            docx_path = f.name
        
        try:
            # Create test DOCX
            doc = Document()
            doc.add_heading("Calculus III Assignment", level=1)
            doc.add_paragraph("Problem 1: Find the derivative of f(x) = x³ + 2x")
            doc.add_paragraph("MATH 53 - Week 4")
            doc.save(docx_path)
            
            # Extract text
            text = await extract_text_from_docx(docx_path)
            
            assert text, "No text extracted from DOCX"
            assert "Calculus" in text or "derivative" in text or "MATH" in text
            assert len(text) > 20
            
            print("✅ DOCX extraction test passed")
            
        finally:
            Path(docx_path).unlink(missing_ok=True)
            
    except ImportError as e:
        pytest.skip(f"Dependencies not installed: {e}")


@pytest.mark.asyncio
async def test_pptx_extraction():
    """Test 4: python-pptx extracts text from PowerPoint slides."""
    try:
        from backend.extraction.pptx import extract_text_from_pptx
        from pptx import Presentation
        
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            pptx_path = f.name
        
        try:
            # Create test PPTX
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            txBox = slide.shapes.add_textbox(10, 10, 500, 100)
            tf = txBox.text_frame
            tf.text = "Physics 51: Modern Physics\nElectromagnetism Chapter\nENGRSCI 45"
            
            prs.save(pptx_path)
            
            # Extract text
            text = await extract_text_from_pptx(pptx_path)
            
            assert text, "No text extracted from PPTX"
            assert "Physics" in text or "Electromagnetism" in text or "ENGR" in text
            assert len(text) > 20
            
            print("✅ PPTX extraction test passed")
            
        finally:
            Path(pptx_path).unlink(missing_ok=True)
            
    except ImportError as e:
        pytest.skip(f"Dependencies not installed: {e}")


@pytest.mark.asyncio
async def test_embedding_stability():
    """Test 5: Embeddings produce stable cosine scores for multilingual text."""
    try:
        from backend.classifier.embeddings import embed, cosine_similarity
        import numpy as np
        
        # Test English text
        text_en = "Introduction to Computer Science and Programming"
        embedding_en = embed([text_en])[0]
        
        assert embedding_en.shape == (384,), f"Wrong embedding shape: {embedding_en.shape}"
        assert abs(np.linalg.norm(embedding_en) - 1.0) < 0.01, "Embedding not normalized"
        
        # Test Filipino text
        text_fil = "Panimula sa Computer Science at Programming"
        embedding_fil = embed([text_fil])[0]
        
        assert embedding_fil.shape == (384,), f"Wrong embedding shape: {embedding_fil.shape}"
        
        # Test similarity
        similarity = cosine_similarity(embedding_en, embedding_fil)
        assert 0.0 <= similarity <= 1.0, f"Similarity out of range: {similarity}"
        assert similarity > 0.5, f"Similar texts should have high similarity: {similarity}"
        
        # Test stability: running twice should give same results
        embedding_en_2 = embed([text_en])[0]
        stability = cosine_similarity(embedding_en, embedding_en_2)
        assert stability > 0.99, f"Embeddings not stable: {stability}"
        
        print("✅ Embedding stability test passed")
        
    except ImportError as e:
        pytest.skip(f"Dependencies not installed: {e}")


@pytest.mark.asyncio
async def test_text_normalization():
    """Test 6: Text normalization handles encoding and whitespace."""
    from backend.extraction.utils import normalize_text, clean_and_truncate
    
    # Test with messy text
    messy = "Hello   \n\n\n   World  \t\t  !!! "
    normalized = normalize_text(messy)
    
    assert "  " not in normalized, "Double spaces not removed"
    assert "\n\n\n" not in normalized, "Multiple newlines not collapsed"
    
    # Test truncation
    long_text = "This is a test. " * 200  # ~3200 chars
    truncated = clean_and_truncate(long_text, 500)
    
    assert len(truncated) <= 510, f"Truncation failed: {len(truncated)}"
    assert truncated[-1] in ['.', '?', ' '], "Truncation should cut at word boundary"
    
    print("✅ Text normalization test passed")


@pytest.mark.asyncio
async def test_course_loading():
    """Test 7: UP Tacloban course seed data loads correctly."""
    try:
        import json
        from pathlib import Path
        
        courses_path = Path("data/up_tacloban_courses.json")
        assert courses_path.exists(), "Course data file not found"
        
        with open(courses_path) as f:
            data = json.load(f)
        
        assert data["institution"] == "University of the Philippines Tacloban College"
        assert len(data["departments"]) == 5
        
        # Count courses
        total_courses = 0
        for dept_name, dept_structure in data["departments"].items():
            for subdept_name, subdept_data in dept_structure.items():
                if isinstance(subdept_data, dict) and 'courses' in subdept_data:
                    total_courses += len(subdept_data['courses'])
        
        assert total_courses > 50, f"Too few courses loaded: {total_courses}"
        
        print(f"✅ Course data test passed ({total_courses} courses)")
        
    except ImportError as e:
        pytest.skip(f"Dependencies not installed: {e}")


# Utility function to run all tests manually
async def run_all_tests():
    """Run all tests when executed as a script."""
    tests = [
        ("Health Endpoint", test_health_endpoint),
        ("PDF Extraction", test_pdf_extraction),
        ("DOCX Extraction", test_docx_extraction),
        ("PPTX Extraction", test_pptx_extraction),
        ("Embedding Stability", test_embedding_stability),
        ("Text Normalization", test_text_normalization),
        ("Course Loading", test_course_loading),
    ]
    
    print("\n🧪 Running Wave 1 Validation Tests\n")
    passed = 0
    failed = 0
    
    for test_name, test_func in tests:
        try:
            print(f"Running: {test_name}...", end=" ", flush=True)
            await test_func()
            passed += 1
        except pytest.skip.Exception as e:
            print(f"⏭️  SKIPPED: {e}")
        except AssertionError as e:
            print(f"❌ FAILED: {e}")
            failed += 1
        except Exception as e:
            print(f"❌ ERROR: {e}")
            failed += 1
    
    print(f"\n📊 Results: {passed} passed, {failed} failed\n")
    return failed == 0


if __name__ == "__main__":
    asyncio.run(run_all_tests())
