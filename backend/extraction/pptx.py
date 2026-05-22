"""PPTX text extraction using python-pptx."""

import logging
from pathlib import Path

from pptx import Presentation

from backend.config import settings
from backend.extraction.utils import clean_and_truncate, normalize_text

logger = logging.getLogger(__name__)


async def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from a PPTX file (PowerPoint presentation).
    
    Args:
        file_path: Path to the PPTX file
        
    Returns:
        Extracted and cleaned text, truncated to max chars
        
    Raises:
        ValueError: If file cannot be read or is corrupted
    """
    try:
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise ValueError(f"PPTX file not found: {file_path}")
        
        if file_path.stat().st_size > settings.max_file_size_mb * 1024 * 1024:
            raise ValueError(
                f"PPTX file too large: {file_path.stat().st_size / 1024 / 1024:.1f}MB "
                f"(max: {settings.max_file_size_mb}MB)"
            )
        
        logger.info(f"Extracting text from PPTX: {file_path}")
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                # Extract from text frames (titles, text boxes)
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
                
                # Extract from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text_parts.append(cell.text)
        
        full_text = "\n".join(text_parts)
        
        if not full_text.strip():
            raise ValueError("No text could be extracted from PPTX")
        
        # Normalize and truncate
        normalized = normalize_text(full_text)
        cleaned = clean_and_truncate(normalized, settings.extraction_max_chars)
        
        logger.info(
            f"PPTX extraction complete: {len(cleaned)} chars from {len(prs.slides)} slides"
        )
        
        return cleaned
        
    except Exception as e:
        logger.error(f"PPTX extraction failed for {file_path}: {e}")
        raise


async def extract_text_from_pptx_with_metadata(file_path: str) -> dict:
    """
    Extract text and metadata from a PPTX file.
    
    Returns dict with keys:
        - text: Extracted text
        - slide_count: Number of slides
        - title: Presentation title (if available)
    """
    try:
        file_path = Path(file_path)
        prs = Presentation(file_path)
        
        # Extract metadata
        core_properties = prs.core_properties
        
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                text_parts.append(cell.text)
        
        full_text = "\n".join(text_parts)
        normalized = normalize_text(full_text)
        cleaned = clean_and_truncate(normalized, settings.extraction_max_chars)
        
        return {
            "text": cleaned,
            "slide_count": len(prs.slides),
            "title": core_properties.title or "",
            "author": core_properties.author or "",
        }
        
    except Exception as e:
        logger.error(f"PPTX metadata extraction failed for {file_path}: {e}")
        raise
